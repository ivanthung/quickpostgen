from pptx import Presentation

prs = Presentation("resources/Metabolic_template.pptx")


print("over")
print (prs.slide_layouts[0].name)
print (prs.slide_layouts[3].name)


text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
print("text_runs: ", text_runs)